#!/usr/bin/env python3
"""Generate PowerPoint presentation for Google Cloud Practice CE Interview.

Enhanced version with:
- Proper introduction and context setting
- About Me slide
- Problem context before solution
- Smooth transitions
- Professional closing

Usage:
    pip install python-pptx
    python generate_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Google Colors
GOOGLE_BLUE = RGBColor(66, 133, 244)
GOOGLE_RED = RGBColor(234, 67, 53)
GOOGLE_YELLOW = RGBColor(251, 188, 5)
GOOGLE_GREEN = RGBColor(52, 168, 83)
DARK_GRAY = RGBColor(60, 64, 67)
LIGHT_GRAY = RGBColor(95, 99, 104)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
CODE_BG = RGBColor(40, 44, 52)
CODE_TEXT = RGBColor(171, 178, 191)


def add_title_slide(prs, title, subtitle="", footer=""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

    # Footer
    if footer:
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.8))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = footer
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

    # Google color bar at bottom
    for i, color in enumerate([GOOGLE_BLUE, GOOGLE_RED, GOOGLE_YELLOW, GOOGLE_GREEN]):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(i * 2.5), Inches(7.1), Inches(2.5), Inches(0.15))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    return slide


def add_content_slide(prs, title, bullets, notes=""):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = GOOGLE_BLUE
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Bullets
    bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.8))
    tf = bullet_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle sub-bullets (indented with spaces) and empty lines
        if bullet == "":
            p.text = ""
            p.space_after = Pt(4)
        elif bullet.startswith("  "):
            p.text = "      " + bullet.strip()
            p.font.size = Pt(16)
            p.level = 1
            p.space_after = Pt(6)
        else:
            p.text = bullet
            p.font.size = Pt(18)
            p.level = 0
            p.space_after = Pt(8)

        p.font.color.rgb = DARK_GRAY

    # Speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide


def add_section_slide(prs, section_number, section_title, subtitle=""):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = GOOGLE_BLUE
    bg.line.fill.background()

    # Section number
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"SECTION {section_number}"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER

    # Section title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(200, 220, 255)
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_code_slide(prs, title, code, notes=""):
    """Add a slide with code block."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = GOOGLE_BLUE
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Code box background
    code_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.3), Inches(9.4), Inches(5.9))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = CODE_BG
    code_bg.line.fill.background()

    # Code text
    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.6))
    tf = code_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(12)
    p.font.name = "Courier New"
    p.font.color.rgb = CODE_TEXT

    # Speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide


def add_comparison_slide(prs, title, left_title, left_content, right_title, right_content, notes=""):
    """Add a two-column comparison slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = GOOGLE_BLUE
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Left column header
    left_header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.3), Inches(4.5), Inches(0.5))
    left_header.fill.solid()
    left_header.fill.fore_color.rgb = GOOGLE_RED
    left_header.line.fill.background()

    left_header_text = slide.shapes.add_textbox(Inches(0.3), Inches(1.35), Inches(4.5), Inches(0.5))
    tf = left_header_text.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Left column content box
    left_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.8), Inches(4.5), Inches(5.4))
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = RGBColor(250, 240, 240)
    left_bg.line.fill.background()

    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.9), Inches(4.3), Inches(5.2))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = left_content
    p.font.size = Pt(11)
    p.font.name = "Courier New"
    p.font.color.rgb = DARK_GRAY

    # Right column header
    right_header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(1.3), Inches(4.5), Inches(0.5))
    right_header.fill.solid()
    right_header.fill.fore_color.rgb = GOOGLE_GREEN
    right_header.line.fill.background()

    right_header_text = slide.shapes.add_textbox(Inches(5.2), Inches(1.35), Inches(4.5), Inches(0.5))
    tf = right_header_text.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Right column content box
    right_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(1.8), Inches(4.5), Inches(5.4))
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = RGBColor(240, 250, 240)
    right_bg.line.fill.background()

    # Right column content
    right_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.9), Inches(4.3), Inches(5.2))
    tf = right_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = right_content
    p.font.size = Pt(11)
    p.font.name = "Courier New"
    p.font.color.rgb = DARK_GRAY

    # Speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide


def add_table_slide(prs, title, headers, rows, notes=""):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = GOOGLE_BLUE
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Table
    num_cols = len(headers)
    num_rows = len(rows) + 1
    row_height = min(0.6, 5.0 / num_rows)

    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.3), Inches(1.4),
        Inches(9.4), Inches(row_height * num_rows)
    ).table

    # Set column widths
    col_width = Inches(9.4 / num_cols)
    for col in table.columns:
        col.width = col_width

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = GOOGLE_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_GRAY

    # Speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide


def create_presentation():
    """Create the full presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ===== SLIDE 1: Title =====
    add_title_slide(
        prs,
        "Unlocking Trapped Data",
        "A Custom RAG Solution for Real Estate Property Valuation",
        "Google Cloud Practice Customer Engineer | Technical Presentation"
    )

    # ===== SLIDE 2: About Me =====
    add_content_slide(
        prs,
        "About Me",
        [
            "[Your Name]",
            "",
            "Background:",
            "  [X] years in cloud architecture and data engineering",
            "  Experience with AI/ML solutions and enterprise systems",
            "  Passion for solving complex data challenges",
            "",
            "Today's Focus:",
            "  Building a custom solution for a real customer problem",
            "  Demonstrating hands-on technical implementation",
            "  Connecting technical decisions to business outcomes",
        ],
    )

    # ===== SLIDE 3: Agenda =====
    add_content_slide(
        prs,
        "Agenda",
        [
            "1. The Challenge",
            "  Understanding the customer's technical blocker (3 min)",
            "",
            "2. The Solution",
            "  Custom RAG architecture with domain-specific cleaning (4 min)",
            "",
            "3. Live Demonstration",
            "  Hands-on walkthrough of the prototype (8 min)",
            "",
            "4. Architectural Decisions",
            "  Trade-offs and design choices (3 min)",
            "",
            "5. Business Value & Next Steps",
            "  ROI, strategic impact, and production roadmap (2 min)",
        ],
    )

    # ===== SLIDE 4: Section - The Challenge =====
    add_section_slide(prs, "1", "The Challenge", "Understanding the Technical Blocker")

    # ===== SLIDE 5: Customer Context =====
    add_content_slide(
        prs,
        "Customer Context",
        [
            "The Customer:",
            "  Large real estate investment firm",
            "  Decades of property valuation data",
            "  Hundreds of analysts performing due diligence daily",
            "",
            "The Situation:",
            "  Ready to adopt Google Cloud",
            "  Want to leverage AI for valuation workflows",
            "  But... they have a critical technical blocker",
            "",
            "The Blocker:",
            "  'Our proprietary data is trapped in siloed, unstructured",
            "   PDF formats that standard RAG tools can't parse correctly'",
        ],
    )

    # ===== SLIDE 6: The Data Reality =====
    add_content_slide(
        prs,
        "The Data Reality: Siloed & Unstructured",
        [
            "Document Types (Silos):",
            "  comps/           - Comparable sales reports",
            "  offering_memo/   - Offering memoranda",
            "  appraisals/      - Appraisal reports",
            "  leases/          - Lease abstracts",
            "  financials/      - NOI statements, rent rolls",
            "",
            "Volume:",
            "  10,000+ documents across silos",
            "  50-200 pages per document",
            "  20+ years of historical data",
            "",
            "Current Pain:",
            "  Analysts spend 2-3 hours finding cap rate evidence",
            "  Inconsistent results across team members",
            "  No audit trail for valuation decisions",
        ],
    )

    # ===== SLIDE 7: Why Standard Tools Fail =====
    add_table_slide(
        prs,
        "Why Standard RAG Tools Fail",
        ["Issue", "Example", "Impact"],
        [
            ["Repeated Headers", '"CONFIDENTIAL REPORT" on every page', "Pollutes search results"],
            ["Page Numbers", '"Page 1 of 2" in content', "Creates meaningless chunks"],
            ["Inconsistent Dates", '"3/7/2025" vs "2025-03-07"', "Queries miss matches"],
            ["Currency Formats", '"$ 1,250,000" vs "$1.25M"', "Retrieval failures"],
            ["Area Units", '"12500 SF" vs "12,500 sq ft"', "Inconsistent indexing"],
            ["Near-Duplicates", "Same paragraph repeated", "Wastes context budget"],
        ],
    )

    # ===== SLIDE 8: Section - The Solution =====
    add_section_slide(prs, "2", "The Solution", "Custom RAG with Domain-Specific Cleaning")

    # ===== SLIDE 9: Solution Overview =====
    add_content_slide(
        prs,
        "The Solution: Custom RAG Architecture",
        [
            "Standard RAG Approach (What Fails):",
            "  PDF -> Extract -> Chunk -> Embed -> Store",
            "  Problem: Noise passes through to the index",
            "",
            "Our Approach (What Works):",
            "  PDF -> Extract -> CUSTOM CLEANING -> Chunk -> Embed -> Store",
            "",
            "The Key Innovation:",
            "  A domain-specific cleaning pipeline that:",
            "  Removes boilerplate before it pollutes the index",
            "  Normalizes formats for consistent matching",
            "  Preserves provenance for auditable citations",
            "  Enables silo-based filtering for targeted queries",
        ],
    )

    # ===== SLIDE 10: Architecture Diagram =====
    add_content_slide(
        prs,
        "End-to-End Architecture",
        [
            "DATA INGESTION PIPELINE:",
            "  [PDF Files] -> [Ingestion] -> [Cleaning] -> [Chunking]",
            "       |             |              |             |",
            "     Input      Extract text    Normalize     300-char",
            "     silos      + metadata      + dedupe      + overlap",
            "",
            "INDEXING & STORAGE:",
            "  [Chunking] -> [Embedding] -> [Vector Store]",
            "       |            |               |",
            "    Chunks      Vectors        ChromaDB",
            "    + meta      (local/API)    (persistent)",
            "",
            "QUERY & RESPONSE:",
            "  [Question] -> [Embed] -> [Search] -> [RAG Engine] -> [Answer]",
            "       |          |           |            |              |",
            "    User       Vector      Top-K       Grounded      Citations",
            "    input      query       chunks      prompt        + score",
        ],
    )

    # ===== SLIDE 11: The Cleaning Pipeline =====
    before_text = """CONFIDENTIAL REPORT
--------------------
Downtown Tower Summary

NOI noted on 3/7/2025
at $ 1,250,000
for 12500 SF
Cap rate is 6.1%

Page 1 of 2
--------------------
CONFIDENTIAL REPORT"""

    after_text = """Downtown Tower Summary

NOI noted on 2025-03-07
at $1,250,000
for 12500 sqft
Cap rate is 6.1%


TRANSFORMATIONS APPLIED:
- Removed repeated header
- Removed page number
- Date: 3/7/2025 -> 2025-03-07
- Currency: $ 1,250,000 -> $1,250,000
- Area: SF -> sqft
- Extracted: cap_rate=6.1%"""

    add_comparison_slide(
        prs,
        "The Cleaning Pipeline: Before & After",
        "BEFORE (Raw PDF)",
        before_text,
        "AFTER (Cleaned)",
        after_text,
    )

    # ===== SLIDE 12: Section - Live Demo =====
    add_section_slide(prs, "3", "Live Demonstration", "Hands-On Technical Walkthrough")

    # ===== SLIDE 13: Demo Overview =====
    add_content_slide(
        prs,
        "Demo Overview",
        [
            "What I'll Demonstrate:",
            "",
            "Step 1: Generate Sample Data",
            "  Create 50 synthetic real estate PDFs across 5 silos",
            "",
            "Step 2: Index the Documents",
            "  Run the full ingestion -> clean -> chunk -> embed pipeline",
            "",
            "Step 3: Query the System",
            "  Ask valuation questions and see grounded answers",
            "",
            "Step 4: Silo Filtering",
            "  Target specific document types for precise results",
            "",
            "Environment: Fully offline, no API keys required",
        ],
    )

    # ===== SLIDE 14: Demo - Generate Data =====
    add_code_slide(
        prs,
        "Demo Step 1: Generate Sample Data",
        """# Generate 50 synthetic real estate PDFs
$ python scripts/generate_50_samples.py

Generated: sample_data_50/comps/comp_001.pdf
Generated: sample_data_50/offering_memo/memo_002.pdf
Generated: sample_data_50/appraisals/appraisal_003.pdf
Generated: sample_data_50/leases/lease_004.pdf
Generated: sample_data_50/financials/financial_005.pdf
...

Generated 50 PDFs in ./sample_data_50

Silo distribution:
  comps:         14 files
  offering_memo: 10 files
  appraisals:    10 files
  leases:         8 files
  financials:     8 files""",
    )

    # ===== SLIDE 15: Demo - Index =====
    add_code_slide(
        prs,
        "Demo Step 2: Index Documents",
        """# Run the full ingestion pipeline
$ re-rag index \\
    --input-dir ./sample_data_50 \\
    --vector-db-path ./vector_db_50 \\
    --collection valuation_50 \\
    --embedding-provider local \\
    --embedding-dimensions 12

# Output:
{
  "status": "ok",
  "documents": 50,           # PDFs ingested
  "clean_segments": 90,      # After cleaning (noise removed!)
  "chunks_indexed": 165,     # After chunking (searchable units)
  "vector_db_path": "./vector_db_50",
  "collection": "valuation_50"
}

# Notice: 50 docs -> 90 segments -> 165 chunks
# The cleaning pipeline removed noise and deduplicated content""",
    )

    # ===== SLIDE 16: Demo - Query =====
    add_code_slide(
        prs,
        "Demo Step 3: Query the System",
        """# Ask a valuation question
$ re-rag ask \\
    --question "What cap rate evidence exists?" \\
    --vector-db-path ./vector_db_50 \\
    --collection valuation_50 \\
    --embedding-provider local \\
    --embedding-dimensions 12 \\
    --llm-provider stub \\
    --top-k 5

ANSWER:
Based on the indexed documents, cap rate evidence includes...

INSUFFICIENT_EVIDENCE: False

CITATIONS:
- chunk_id=abc123... doc_id=def456... page_span=1-1 score=0.847
- chunk_id=ghi789... doc_id=jkl012... page_span=2-2 score=0.723
- chunk_id=mno345... doc_id=pqr678... page_span=1-1 score=0.651""",
    )

    # ===== SLIDE 17: Demo - Silo Filtering =====
    add_code_slide(
        prs,
        "Demo Step 4: Silo-Based Filtering",
        """# Filter queries to specific document types
$ re-rag ask \\
    --question "Show me lease terms and rent escalations" \\
    --silo leases \\                    # <-- Only search leases/
    --vector-db-path ./vector_db_50 \\
    --collection valuation_50 \\
    --embedding-provider local \\
    --embedding-dimensions 12 \\
    --llm-provider stub \\
    --top-k 3

ANSWER:
Based on lease abstracts, typical terms include...

CITATIONS:  (All from leases/ silo)
- chunk_id=... doc_id=... page_span=1-1 score=0.409 (lease_008.pdf)
- chunk_id=... doc_id=... page_span=1-1 score=0.078 (lease_021.pdf)
- chunk_id=... doc_id=... page_span=1-1 score=0.074 (lease_043.pdf)""",
    )

    # ===== SLIDE 18: Section - Architecture =====
    add_section_slide(prs, "4", "Architectural Decisions", "Trade-offs and Design Choices")

    # ===== SLIDE 19: Technology Stack =====
    add_table_slide(
        prs,
        "Technology Stack",
        ["Component", "Choice", "Why"],
        [
            ["Language", "Python 3.10+", "Fast iteration, rich ML ecosystem"],
            ["PDF Extraction", "pypdf", "Reliable page-level extraction"],
            ["Vector DB", "ChromaDB", "Persistent, local-first, easy migration"],
            ["Embeddings", "Pluggable adapter", "Local for dev, Vertex AI for prod"],
            ["LLM", "Pluggable adapter", "Stub for demo, Gemini for prod"],
            ["Testing", "pytest (30 tests)", "Full pipeline coverage"],
        ],
    )

    # ===== SLIDE 20: Key Design Decisions =====
    add_content_slide(
        prs,
        "Key Design Decisions",
        [
            "Decision 1: Why ChromaDB?",
            "  + Local-first: No cloud dependency for prototype",
            "  + Persistent: Survives restarts, real database behavior",
            "  + Migration path: Easy swap to Vertex AI Vector Search",
            "",
            "Decision 2: Why Character-Based Chunking?",
            "  + Deterministic: Same input = same output, every time",
            "  + Model-agnostic: No tokenizer dependency",
            "  + Clear boundaries: Precise citation page spans",
            "",
            "Decision 3: Why Pluggable Adapters?",
            "  + Demo reliability: No API failures during presentation",
            "  + Environment flexibility: Local dev, cloud production",
            "  + Testing: Deterministic CI/CD pipelines",
        ],
    )

    # ===== SLIDE 21: Safety & Trust =====
    add_content_slide(
        prs,
        "Safety & Trust Features",
        [
            "Challenge: LLMs can hallucinate, especially with domain data",
            "",
            "Safeguard 1: Grounded Prompts",
            "  LLM instructed to use ONLY provided context",
            "  Explicit instruction: 'If evidence is incomplete, say so'",
            "",
            "Safeguard 2: Citation Tracking",
            "  Every answer includes: chunk_id, doc_id, page_span, score",
            "  Full audit trail back to source documents",
            "",
            "Safeguard 3: No-Evidence Handling",
            "  Returns insufficient_evidence=True when uncertain",
            "  Never fabricates data - asks for more documents instead",
            "",
            "Safeguard 4: Configurable Thresholds",
            "  Minimum similarity score, context budget limits",
        ],
    )

    # ===== SLIDE 22: Section - Business Value =====
    add_section_slide(prs, "5", "Business Value", "ROI, Strategic Impact, and Next Steps")

    # ===== SLIDE 23: Quantified ROI =====
    add_table_slide(
        prs,
        "Business Impact: Quantified ROI",
        ["Metric", "Before", "After", "Impact"],
        [
            ["Time to find evidence", "2-3 hours", "30 seconds", "99% reduction"],
            ["Queries per analyst/day", "2-3", "50+", "20x throughput"],
            ["Answer consistency", "Variable", "Standardized", "Risk reduction"],
            ["Audit trail", "Manual notes", "Automatic citations", "Compliance ready"],
            ["Onboarding time", "Weeks", "Days", "Faster ramp-up"],
        ],
    )

    # ===== SLIDE 24: Strategic Value =====
    add_content_slide(
        prs,
        "Strategic Value",
        [
            "1. Unlock Trapped Data",
            "  Decades of institutional knowledge becomes searchable",
            "  Competitive intelligence at analysts' fingertips",
            "",
            "2. Accelerate Deal Cycles",
            "  Faster due diligence = more deals closed",
            "  Consistent analysis across team members",
            "",
            "3. Reduce Risk",
            "  Auditable citations reduce unsupported claims",
            "  No hallucinated valuations",
            "",
            "4. Scalable Pattern",
            "  Architecture applies to other document-heavy workflows",
            "  Legal, compliance, research departments",
        ],
    )

    # ===== SLIDE 25: Production Roadmap =====
    add_content_slide(
        prs,
        "Production Roadmap",
        [
            "Phase 1: Pilot (Weeks 1-6)",
            "  Deploy to Cloud Run",
            "  Integrate Cloud Storage + Document AI",
            "  Add real customer documents",
            "",
            "Phase 2: Production (Weeks 6-12)",
            "  Vertex AI embeddings + Gemini LLM",
            "  Vertex AI Vector Search or AlloyDB",
            "  SLOs, monitoring, alerting",
            "",
            "Phase 3: Scale (Weeks 12+)",
            "  VPC Service Controls",
            "  Multi-zone DR",
            "  Continuous quality regression",
            "",
            "Success Criteria:",
            "  95% relevant citations | p95 latency < 3s | Zero hallucinations",
        ],
    )

    # ===== SLIDE 26: Summary =====
    add_content_slide(
        prs,
        "Summary: Key Takeaways",
        [
            "The Problem",
            "  Proprietary valuation data trapped in unstructured PDFs",
            "",
            "The Solution",
            "  Custom RAG with domain-specific cleaning pipeline",
            "",
            "The Differentiator",
            "  Noise removal + normalization BEFORE indexing",
            "",
            "The Result",
            "  Answers in seconds, not hours",
            "  Full citation traceability",
            "  Safe handling of insufficient evidence",
            "",
            "The Path Forward",
            "  Clear roadmap to Google Cloud production",
        ],
    )

    # ===== SLIDE 27: Thank You =====
    add_title_slide(
        prs,
        "Thank You",
        "Questions & Discussion",
        "[Your Name] | [Your Email]"
    )

    # ===== SLIDE 28: Appendix - Q&A Prep =====
    add_content_slide(
        prs,
        "Appendix: Anticipated Q&A",
        [
            "Technical Questions:",
            "  Q: Why not LangChain/LlamaIndex?",
            "  A: They don't handle real estate noise patterns",
            "",
            "  Q: How do you handle scanned PDFs?",
            "  A: Flag as scan_suspected, route to Document AI for OCR",
            "",
            "  Q: Why local embeddings?",
            "  A: Demo reliability; same architecture supports Vertex AI",
            "",
            "Business Questions:",
            "  Q: What's the timeline?",
            "  A: 6 weeks pilot, 12 weeks production",
            "",
            "  Q: How do you ensure accuracy?",
            "  A: Citations link every claim to source; no-evidence refusal",
        ],
    )

    # Save presentation
    output_path = "Real_Estate_RAG_Interview_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()